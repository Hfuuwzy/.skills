import argparse
import sys
from pathlib import Path

from pptx import Presentation


def iter_text_runs(path: Path):
    prs = Presentation(path)
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            shape_text = shape.text_frame.text.strip().replace("\n", " / ")
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    size = run.font.size.pt if run.font.size else None
                    yield slide_index, size, run.text.strip(), shape_text


def main() -> int:
    parser = argparse.ArgumentParser(description="Fail if any visible PPTX text run is below a font-size floor.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--min-size", type=float, default=16.0)
    args = parser.parse_args()

    violations = []
    missing = []
    for slide, size, run_text, shape_text in iter_text_runs(args.pptx):
        if size is None:
            missing.append((slide, run_text, shape_text))
        elif size < args.min_size:
            violations.append((slide, size, run_text, shape_text))

    if violations or missing:
        print(f"FAIL: font floor {args.min_size:g}pt not met in {args.pptx}")
        for slide, size, run_text, shape_text in violations[:80]:
            print(f"slide {slide:02d}: {size:g}pt run={run_text!r} shape={shape_text[:120]!r}")
        if len(violations) > 80:
            print(f"... {len(violations) - 80} more below-floor runs")
        for slide, run_text, shape_text in missing[:20]:
            print(f"slide {slide:02d}: missing explicit size run={run_text!r} shape={shape_text[:120]!r}")
        if len(missing) > 20:
            print(f"... {len(missing) - 20} more runs without explicit size")
        return 1

    print(f"PASS: all visible text runs are >= {args.min_size:g}pt")
    return 0


if __name__ == "__main__":
    sys.exit(main())
